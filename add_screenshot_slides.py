"""Add new screenshot slides to IPCMS_v2.pptx with images from ppt_imagwes/.

Each of the 5 screenshots gets its own slide with a caption.
Inserts before the THANK YOU slide and updates page numbers throughout.
"""

import os
from copy import deepcopy
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

SRC = "IPCMS_v2.pptx"
IMG_DIR = r"E:\Integrated_Patient_Care_Management_System\ppt_imagwes"

TEAL = RGBColor(0x14, 0xB8, 0xA6)
DEEP = RGBColor(0x0E, 0x4F, 0x8A)
SLATE = RGBColor(0x0F, 0x17, 0x2A)
MUTED = RGBColor(0x64, 0x74, 0x8B)
BODY = RGBColor(0x33, 0x4E, 0x68)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG_TINT = RGBColor(0xF1, 0xF5, 0xF9)
CARD_FILL = RGBColor(0xF8, 0xFA, 0xFC)
DIVIDER = RGBColor(0xCB, 0xD5, 0xE1)


def in_to_emu(x: float) -> int:
    return int(round(x * 914400))


def add_header(prs, slide, eyebrow, title, subtitle):
    """Eyebrow + big title + subtitle + horizontal rule."""
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    in_to_emu(0.55), in_to_emu(0.55),
                                    in_to_emu(0.12), in_to_emu(0.55))
    accent.line.fill.background()
    accent.fill.solid()
    accent.fill.fore_color.rgb = TEAL

    tb = slide.shapes.add_textbox(in_to_emu(0.78), in_to_emu(0.5),
                                   in_to_emu(12.0), in_to_emu(0.28))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = eyebrow
    r.font.name = "Calibri"
    r.font.size = Pt(10)
    r.font.bold = True
    r.font.color.rgb = TEAL

    tb = slide.shapes.add_textbox(in_to_emu(0.78), in_to_emu(0.78),
                                   in_to_emu(12.2), in_to_emu(0.55))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = "Calibri"
    r.font.size = Pt(26)
    r.font.bold = True
    r.font.color.rgb = SLATE

    tb = slide.shapes.add_textbox(in_to_emu(0.78), in_to_emu(1.35),
                                   in_to_emu(12.2), in_to_emu(0.4))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = subtitle
    r.font.name = "Calibri"
    r.font.size = Pt(12)
    r.font.color.rgb = MUTED

    rule = slide.shapes.add_connector(1, in_to_emu(0.78), in_to_emu(1.78),
                                      in_to_emu(12.6), in_to_emu(1.78))
    rule.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
    rule.line.width = Pt(0.75)


def add_footer(prs, slide, page_num, total):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, in_to_emu(7.18),
                                  prs.slide_width, in_to_emu(0.32))
    bar.line.fill.background()
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(0xF8, 0xFA, 0xFC)

    tb = slide.shapes.add_textbox(in_to_emu(0.4), in_to_emu(7.22),
                                   in_to_emu(9.0), in_to_emu(0.25))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "IPCMS — Integrated Patient Care Management System"
    r.font.name = "Calibri"
    r.font.size = Pt(9)
    r.font.color.rgb = MUTED

    tb = slide.shapes.add_textbox(in_to_emu(11.6), in_to_emu(7.22),
                                   in_to_emu(1.4), in_to_emu(0.25))
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = f"{page_num} / {total}"
    r.font.name = "Calibri"
    r.font.size = Pt(9)
    r.font.color.rgb = MUTED


def add_text(slide, left, top, width, height, text,
             *, size=11, bold=False, color=BODY, name="Calibri",
             align=None):
    tb = slide.shapes.add_textbox(in_to_emu(left), in_to_emu(top),
                                   in_to_emu(width), in_to_emu(height))
    tf = tb.text_frame
    tf.word_wrap = True
    if isinstance(text, str):
        lines = [text]
    else:
        lines = text
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        if align is not None:
            p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.name = name
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return tb


def add_card(slide, left, top, width, height, *,
             fill=CARD_FILL, line_rgb=RGBColor(0xE2, 0xE8, 0xF0)):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   in_to_emu(left), in_to_emu(top),
                                   in_to_emu(width), in_to_emu(height))
    card.fill.solid()
    card.fill.fore_color.rgb = fill
    card.line.color.rgb = line_rgb
    card.line.width = Pt(0.75)
    card.shadow.inherit = False
    return card


def add_image_with_caption(prs, slide, image_path, eyebrow, title, caption):
    """Add a screenshot slide with the image and a caption box."""
    add_header(prs, slide, eyebrow, title, caption)

    # Image card background
    add_card(slide, 0.55, 1.95, 12.2, 4.4,
             fill=WHITE, line_rgb=DIVIDER)

    # Place image, centered, fit within the card area
    # Card is 12.2 wide x 4.4 tall starting at (0.55, 1.95)
    # Image area: x=0.7, y=2.1, max w=11.9, max h=4.1
    try:
        from PIL import Image
        with Image.open(image_path) as im:
            iw, ih = im.size
        aspect = iw / ih
        max_w = 11.9
        max_h = 4.1
        if aspect >= (max_w / max_h):
            # width-constrained
            disp_w = max_w
            disp_h = max_w / aspect
        else:
            # height-constrained
            disp_h = max_h
            disp_w = max_h * aspect
        # center within card
        cx = 0.55 + (12.2 - disp_w) / 2
        cy = 1.95 + (4.4 - disp_h) / 2
        slide.shapes.add_picture(image_path,
                                  in_to_emu(cx), in_to_emu(cy),
                                  width=in_to_emu(disp_w),
                                  height=in_to_emu(disp_h))
    except Exception as e:
        # Fallback: place at fixed position
        slide.shapes.add_picture(image_path,
                                  in_to_emu(0.7), in_to_emu(2.1),
                                  width=in_to_emu(11.9),
                                  height=in_to_emu(4.1))

    # Caption pill at the bottom of the card
    add_card(slide, 0.55, 6.45, 12.2, 0.65,
             fill=BG_TINT, line_rgb=DIVIDER)
    add_text(slide, 0.75, 6.5, 11.8, 0.55,
             caption, size=11, color=BODY)


def move_slide(prs, old_index, new_index):
    sldIdLst = prs.slides._sldIdLst
    slides = list(sldIdLst)
    sldIdLst.remove(slides[old_index])
    sldIdLst.insert(new_index, slides[old_index])


def slide_first_text(slide):
    for s in slide.shapes:
        if s.has_text_frame and s.text_frame.text.strip():
            return s.text_frame.text.strip()
    return ""


def purge_existing_screenshot_slides(prs):
    """Remove any pre-existing screenshot slides (idempotency)."""
    targets = {"SECTION 13", "SCREENSHOT 01", "SCREENSHOT 02",
               "SCREENSHOT 03", "SCREENSHOT 04", "SCREENSHOT 05"}
    for idx in range(len(prs.slides) - 1, -1, -1):
        first = slide_first_text(prs.slides[idx]).split("\n")[0].strip()
        if first in targets:
            remove_slide(prs, idx)


def remove_slide(prs, index):
    sldIdLst = prs.slides._sldIdLst
    slides = list(sldIdLst)
    sldId = slides[index]
    rId = sldId.get(qn("r:id"))
    prs.part.drop_rel(rId)
    sldIdLst.remove(sldId)


def main():
    prs = Presentation(SRC)
    purge_existing_screenshot_slides(prs)

    # Collect images from the folder
    images = sorted([
        f for f in os.listdir(IMG_DIR)
        if f.lower().endswith((".png", ".jpg", ".jpeg"))
    ])

    if not images:
        print(f"No images found in {IMG_DIR}")
        return

    print(f"Found {len(images)} images: {images}")

    # Captions for each screenshot (based on filename timestamps and order)
    captions = [
        ("SCREENSHOT 01",
         "Application login interface",
         "Polished email + password and Google OAuth sign-in flow with role detection and the glassmorphism theme."),
        ("SCREENSHOT 02",
         "Patient dashboard",
         "Patient view of appointments, history, prescriptions, and the AI care assistant."),
        ("SCREENSHOT 03",
         "Doctor consultation view",
         "Doctor-side schedule, AI-generated pre-visit brief, vitals, notes, and prescription writer."),
        ("SCREENSHOT 04",
         "Pharmacy & medicine catalogue",
         "Patient-facing pharmacy with browse, buy and book flows, powered by the digital medicine catalogue."),
        ("SCREENSHOT 05",
         "Admin analytics dashboard",
         "Admin analytics: visits, revenue, doctor load — backed by pandas DataFrames from the same MySQL DB."),
    ]

    # Fallback if more images than captions
    while len(captions) < len(images):
        i = len(captions) + 1
        captions.append((
            f"SCREENSHOT {i:02d}",
            f"Application screenshot {i}",
            f"Capture from the running application — section {i:02d} of the demo walkthrough."))

    # Record how many slides existed before
    existing = len(prs.slides)
    # We are inserting len(images) new screenshot slides + 1 section header.
    # Section header is optional; keep it simple: just the screenshot slides.
    n_new = len(images)
    final_total = existing + n_new

    # The THANK YOU slide is the last slide currently (index existing-1).
    thank_you_idx = existing - 1

    # Build the new screenshot slides. They'll be appended to the end,
    # then we'll move them so they sit before THANK YOU.
    new_slide_indices = []
    for i, (img_name, (eyebrow, title, caption)) in enumerate(zip(images, captions)):
        img_path = os.path.join(IMG_DIR, img_name)
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
        # Page number for this slide will be assigned after we know final_total
        # but at this point final_total is known (we appended them).
        # We re-assign after moving.
        add_image_with_caption(prs, slide, img_path, eyebrow, title, caption)
        new_slide_indices.append(len(prs.slides) - 1)

    # Move the new slides (currently at the end) to be just before THANK YOU.
    # We do this in reverse order so earlier indices remain valid.
    # The new slides are currently at indices [existing, existing+1, ..., existing+n_new-1].
    # THANK YOU is at index existing-1.
    # We want the new slides to occupy indices [existing-1, existing, ..., existing+n_new-2].
    # Insert in reverse: last new slide goes to index existing-1, then previous goes to existing-1, etc.
    target = existing - 1
    for idx in reversed(new_slide_indices):
        move_slide(prs, idx, target)

    # Update footer page numbers throughout the deck.
    total_slides = len(prs.slides)
    for idx in range(total_slides):
        slide = prs.slides[idx]
        for shp in slide.shapes:
            if shp.has_text_frame:
                for p in shp.text_frame.paragraphs:
                    for r in p.runs:
                        if r.text and ("/" in r.text):
                            txt = r.text.strip()
                            # match patterns like "12 / 19", "1 / 14"
                            import re
                            if re.match(r"^\d+\s*/\s*\d+$", txt):
                                left, _ = txt.split("/")
                                r.text = f"{left.strip()} / {final_total}"

    prs.save(SRC)
    print(f"Saved {SRC} with {final_total} slides (added {n_new} screenshot slides).")


if __name__ == "__main__":
    main()
