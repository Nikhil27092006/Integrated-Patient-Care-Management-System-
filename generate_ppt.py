"""
Generate PCMHS Project Presentation (PowerPoint)
A comprehensive, human-curated PPT covering definition, working, components,
AI implementations, difficulties, features, frontend/backend and tech stack.

Layout rules used everywhere:
 - Title block lives between y=0.4 and y=1.85
 - Footer lives between y=7.18 and y=7.5
 - All content cards live between y=2.0 and y=7.0  (~5.0 inches of usable height)
 - Text font sizes are conservative (10–13pt) to avoid cutoffs
 - Cards are sized so 2x3 grids fit the page without overlap
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


# ── Color palette ─────────────────────────────────────────────────
PRIMARY     = RGBColor(0x0E, 0x4F, 0x8A)   # deep medical blue
ACCENT      = RGBColor(0x14, 0xB8, 0xA6)   # teal accent
SECONDARY   = RGBColor(0x06, 0x6F, 0xA8)   # sky blue
PURPLE      = RGBColor(0x7C, 0x3A, 0xED)
PINK        = RGBColor(0xDB, 0x27, 0x77)
DARK_TEXT   = RGBColor(0x0F, 0x17, 0x2A)
BODY_TEXT   = RGBColor(0x33, 0x4E, 0x68)
MUTED_TEXT  = RGBColor(0x64, 0x74, 0x8B)
LIGHT_BG    = RGBColor(0xF1, 0xF5, 0xF9)
SOFT_BG     = RGBColor(0xE0, 0xF2, 0xFE)
CARD_BG     = RGBColor(0xFF, 0xFF, 0xFF)
PAGE_BG     = RGBColor(0xFB, 0xFC, 0xFD)
ACCENT_SOFT = RGBColor(0xCC, 0xFB, 0xF1)
WARN        = RGBColor(0xEA, 0x58, 0x0C)
DIVIDER     = RGBColor(0xCB, 0xD5, 0xE1)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)


# ── Presentation setup ─────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


# ── Helpers ────────────────────────────────────────────────────────
def add_rect(slide, x, y, w, h, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    return shp


def add_round_rect(slide, x, y, w, h, fill, line=None, radius=0.05):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.adjustments[0] = radius
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    return shp


def add_line(slide, x1, y1, x2, y2, color, weight=1.5):
    ln = slide.shapes.add_connector(1, x1, y1, x2, y2)
    ln.line.color.rgb = color
    ln.line.width = Pt(weight)
    return ln


def add_text(slide, x, y, w, h, text, *,
             size=12, bold=False, color=DARK_TEXT, italic=False,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             font="Calibri", spacing=1.2):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(36000)
    tf.margin_top = tf.margin_bottom = Emu(18000)
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
    return tb


def add_bullets(slide, x, y, w, h, items, *,
                size=11, color=BODY_TEXT, bullet="•",
                spacing=1.25, bold_first=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(36000)
    tf.margin_top = tf.margin_bottom = Emu(18000)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = spacing
        p.space_before = Pt(3)
        p.space_after = Pt(2)
        r = p.add_run()
        r.text = f"{bullet}  {item}"
        r.font.name = "Calibri"
        r.font.size = Pt(size)
        r.font.color.rgb = color
        if bold_first and i == 0:
            r.font.bold = True
    return tb


def page_bg(slide):
    add_rect(slide, 0, 0, SW, SH, PAGE_BG)


def add_footer(slide, page, total):
    add_rect(slide, 0, Inches(7.18), SW, Inches(0.32), LIGHT_BG)
    add_text(slide, Inches(0.4), Inches(7.22), Inches(9), Inches(0.25),
             "PCMHS — Patient Care Management System for Healthcare Services",
             size=9, color=MUTED_TEXT, italic=True)
    add_text(slide, Inches(11.6), Inches(7.22), Inches(1.4), Inches(0.25),
             f"{page} / {total}", size=9, color=MUTED_TEXT,
             align=PP_ALIGN.RIGHT)


def add_title(slide, kicker, title, sub=None):
    """Consistent title block: kicker + title + optional sub, all between 0.4 and 1.85."""
    add_rect(slide, Inches(0.55), Inches(0.55),
             Inches(0.12), Inches(0.55), ACCENT)
    add_text(slide, Inches(0.78), Inches(0.50),
             Inches(12), Inches(0.28),
             kicker.upper(), size=10, bold=True, color=ACCENT)
    add_text(slide, Inches(0.78), Inches(0.78),
             Inches(12.2), Inches(0.55),
             title, size=26, bold=True, color=DARK_TEXT)
    if sub:
        add_text(slide, Inches(0.78), Inches(1.35),
                 Inches(12.2), Inches(0.40),
                 sub, size=12, color=MUTED_TEXT, italic=True)
    add_line(slide, Inches(0.78), Inches(1.78),
             Inches(12.6), Inches(1.78), DIVIDER, 0.75)


# ────────────────────────────────────────────────────────────────────
# SLIDE 1 — Cover (improved, more compact)
# ────────────────────────────────────────────────────────────────────
def slide_1():
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SW, SH, LIGHT_BG)

    # left panel
    add_rect(s, 0, 0, Inches(5.6), SH, PRIMARY)
    add_rect(s, Inches(-1.4), Inches(5.2), Inches(3.2), Inches(3.2), ACCENT)
    add_rect(s, Inches(3.4),  Inches(-1.2), Inches(2.6), Inches(2.6), SECONDARY)

    # top pill
    add_round_rect(s, Inches(0.6), Inches(0.5), Inches(2.5), Inches(0.38),
                   ACCENT, radius=0.5)
    add_text(s, Inches(0.6), Inches(0.5), Inches(2.5), Inches(0.38),
             "AI  •  HEALTHCARE  •  PYTHON",
             size=9.5, bold=True, color=DARK_TEXT,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # mini stat row
    add_text(s, Inches(0.6), Inches(1.05), Inches(2.5), Inches(0.3),
             "PROJECT SHOWCASE", size=11, bold=True, color=ACCENT)
    add_text(s, Inches(3.3), Inches(1.05), Inches(2.0), Inches(0.3),
             "v1.0 · 2026", size=11, italic=True, color=ACCENT_SOFT,
             align=PP_ALIGN.RIGHT)

    # main title
    add_text(s, Inches(0.6), Inches(1.5), Inches(5), Inches(0.5),
             "Integrated Patient", size=32, bold=True, color=WHITE,
             spacing=1.0)
    add_text(s, Inches(0.6), Inches(2.15), Inches(5), Inches(0.5),
             "Care Management", size=32, bold=True, color=WHITE,
             spacing=1.0)
    add_text(s, Inches(0.6), Inches(2.8), Inches(5), Inches(0.5),
             "System", size=32, bold=True, color=ACCENT,
             spacing=1.0)

    # accent underline
    add_line(s, Inches(0.6), Inches(3.55), Inches(3.6), Inches(3.55), ACCENT, 2.5)

    # tagline
    add_text(s, Inches(0.6), Inches(3.7), Inches(5), Inches(0.4),
             "PCMHS · Where AI meets everyday clinical care",
             size=12, italic=True, color=ACCENT_SOFT)

    # 3 KPI tiles — compact, fit between tagline and footer
    kpis = [
        ("3",  "role-aware\ndashboards"),
        ("4",  "AI features\nbuilt in"),
        ("20+", "libraries\ndocumented"),
    ]
    kx0, ky = Inches(0.6), Inches(4.5)
    kw, kh = Inches(1.55), Inches(1.15)
    kgx = Inches(0.15)
    for i, (num, lbl) in enumerate(kpis):
        x = kx0 + (kw + kgx) * i
        add_round_rect(s, x, ky, kw, kh, ACCENT, radius=0.12)
        add_text(s, x, ky + Inches(0.1), kw, Inches(0.55),
                 num, size=22, bold=True, color=DARK_TEXT,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x, ky + Inches(0.62), kw, Inches(0.5),
                 lbl, size=9, color=DARK_TEXT,
                 align=PP_ALIGN.CENTER, spacing=1.1)

    # built-by block
    add_text(s, Inches(0.6), Inches(6.05), Inches(5), Inches(0.25),
             "BUILT BY", size=9, color=ACCENT_SOFT, bold=True)
    add_text(s, Inches(0.6), Inches(6.3), Inches(5), Inches(0.3),
             "Nikhil  •  Aug 2026", size=13, color=WHITE, bold=True)
    add_text(s, Inches(0.6), Inches(6.65), Inches(5), Inches(0.3),
             "Streamlit · LangChain · Groq · Gemini · MySQL",
             size=9.5, italic=True, color=ACCENT_SOFT)

    # right card
    add_rect(s, Inches(6.0), Inches(0.5), Inches(7), Inches(6.5),
             CARD_BG, line=DIVIDER)
    add_rect(s, Inches(6.0), Inches(0.5), Inches(7), Inches(0.14), ACCENT)

    # header
    add_text(s, Inches(6.3), Inches(0.85), Inches(6.5), Inches(0.4),
             "Inside this deck", size=18, bold=True, color=DARK_TEXT)
    add_text(s, Inches(6.3), Inches(1.25), Inches(6.5), Inches(0.5),
             "14 slides — definitions, working, AI layers, "
             "the full library list, features, and the rough edges.",
             size=10.5, color=MUTED_TEXT, italic=True, spacing=1.25)

    # "14 SLIDES" badge top-right of card
    add_round_rect(s, Inches(11.8), Inches(0.85), Inches(1.0), Inches(0.4),
                   ACCENT, radius=0.5)
    add_text(s, Inches(11.8), Inches(0.85), Inches(1.0), Inches(0.4),
             "14 SLIDES", size=10, bold=True, color=DARK_TEXT,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # 6 chips as 3x2 grid (3 columns, 2 rows) — fits more compactly
    chips = [
        ("What it is",         PRIMARY),
        ("How it works",       SECONDARY),
        ("Components",         ACCENT),
        ("AI implementations", PURPLE),
        ("Stack & libraries",  PINK),
        ("Difficulties",       WARN),
    ]
    cx0, cy0 = Inches(6.3), Inches(2.0)
    cw, ch = Inches(2.05), Inches(1.15)
    gx, gy = Inches(0.15), Inches(0.2)
    for i, (lbl, col) in enumerate(chips):
        r, c = i // 3, i % 3
        x = cx0 + (cw + gx) * c
        y = cy0 + (ch + gy) * r
        add_round_rect(s, x, y, cw, ch, SOFT_BG, line=col, radius=0.15)
        # small number tag
        add_round_rect(s, x + Inches(0.15), y + Inches(0.15),
                       Inches(0.4), Inches(0.4),
                       col, radius=0.5)
        add_text(s, x + Inches(0.15), y + Inches(0.15),
                 Inches(0.4), Inches(0.4),
                 f"0{i+1}", size=9, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # label centred to the right of the badge
        add_text(s, x + Inches(0.65), y + Inches(0.15),
                 cw - Inches(0.7), Inches(0.85),
                 lbl, size=11, bold=True, color=col,
                 anchor=MSO_ANCHOR.MIDDLE)

    # bottom signature
    add_line(s, Inches(6.3), Inches(5.95), Inches(12.6), Inches(5.95),
             DIVIDER, 0.75)
    add_text(s, Inches(6.3), Inches(6.05), Inches(6.5), Inches(0.3),
             "Curated by hand · Drafted with AI · Reviewed by the author",
             size=10, italic=True, color=MUTED_TEXT)
    add_text(s, Inches(6.3), Inches(6.4), Inches(6.5), Inches(0.3),
             "For 3rd-year / final-year project showcase",
             size=9.5, color=MUTED_TEXT)
    add_text(s, Inches(6.3), Inches(6.7), Inches(6.5), Inches(0.3),
             "swipe →", size=10, color=ACCENT, bold=True)


# ────────────────────────────────────────────────────────────────────
# SLIDE 2 — Definition
# ────────────────────────────────────────────────────────────────────
def slide_2():
    s = prs.slides.add_slide(BLANK)
    page_bg(s)
    add_title(s, "Section 01", "What is PCMHS?",
             "A working definition, the way I'd explain it to a friend.")

    # left definition card
    add_round_rect(s, Inches(0.55), Inches(2.0), Inches(6.0), Inches(4.95),
                   CARD_BG, line=DIVIDER, radius=0.04)
    add_rect(s, Inches(0.55), Inches(2.0), Inches(0.18), Inches(4.95), ACCENT)

    add_text(s, Inches(0.95), Inches(2.15), Inches(5.4), Inches(0.4),
             "IN ONE SENTENCE", size=11, bold=True, color=ACCENT)
    add_text(s, Inches(0.95), Inches(2.55), Inches(5.4), Inches(2.0),
             "PCMHS is a web-based hospital workflow app that lets "
             "patients, doctors and admins work off the same data — "
             "with an AI layer that helps everyone make better, "
             "faster decisions.",
             size=14, color=DARK_TEXT, spacing=1.4)

    add_line(s, Inches(0.95), Inches(4.7), Inches(6.35), Inches(4.7),
             DIVIDER, 0.5)
    add_text(s, Inches(0.95), Inches(4.85), Inches(5.4), Inches(0.4),
             "LONGER VERSION", size=11, bold=True, color=PRIMARY)
    add_bullets(s, Inches(0.95), Inches(5.25), Inches(5.4), Inches(1.6), [
        "End-to-end OPD/IPD flow: register, book, consult, prescribe, follow-up.",
        "Three role-aware dashboards — what a patient sees is not what a doctor sees.",
        "Two AI assistants — a floating voice chatbot and a database-aware care copilot.",
    ], size=11)

    # right — what problem
    add_text(s, Inches(7.0), Inches(2.0), Inches(5.9), Inches(0.4),
             "What problem it tries to solve",
             size=14, bold=True, color=PRIMARY)

    items = [
        ("Scattered records", "Patient history lives in different places "
                              "for every doctor."),
        ("Manual triage",     "Reception reads forms and routes people by hand."),
        ("Generic advice",    "Patients Google symptoms; doctors can't scale 1-to-1 chat."),
        ("Slow onboarding",   "Doctors, admins and patients each need a "
                              "fast, role-specific start."),
    ]
    y = Inches(2.5)
    for i, (h, b) in enumerate(items):
        add_round_rect(s, Inches(7.0), y, Inches(5.9), Inches(1.0),
                       CARD_BG, line=DIVIDER, radius=0.06)
        add_round_rect(s, Inches(7.0), y, Inches(0.5), Inches(1.0),
                       ACCENT, radius=0.06)
        add_text(s, Inches(7.0), y, Inches(0.5), Inches(1.0),
                 f"0{i+1}", size=13, bold=True, color=DARK_TEXT,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(7.65), y + Inches(0.12), Inches(5.1), Inches(0.35),
                 h, size=12, bold=True, color=DARK_TEXT)
        add_text(s, Inches(7.65), y + Inches(0.48), Inches(5.1), Inches(0.5),
                 b, size=10, color=MUTED_TEXT, italic=True, spacing=1.2)
        y += Inches(1.08)


# ────────────────────────────────────────────────────────────────────
# SLIDE 3 — How it works
# ────────────────────────────────────────────────────────────────────
def slide_3():
    s = prs.slides.add_slide(BLANK)
    page_bg(s)
    add_title(s, "Section 02", "How it works",
             "The journey of a single appointment, end to end.")

    steps = [
        ("Sign in",        "Email + password OR Google OAuth.\nRole detected from DB."),
        ("Pick a doctor",  "Patient browses specialties,\nslots and picks a time."),
        ("AI pre-check",   "Care copilot reads the patient's\nhistory + symptom prompt."),
        ("Consultation",   "Doctor reviews AI summary,\nnotes, prescribes meds."),
        ("Follow-up",      "Email receipt + chat ref.\nVoice chatbot stays 24×7."),
    ]
    n = len(steps)
    content_left  = Inches(0.55)
    content_right = Inches(12.75)
    total_w = content_right - content_left                 # 12.20"
    card_w = Inches(2.10)
    gap    = (total_w - card_w * n) / (n - 1)
    cy = Inches(2.3)

    for i, (title, body) in enumerate(steps):
        x = content_left + (card_w + gap) * i
        add_round_rect(s, x, cy, card_w, Inches(2.3),
                       CARD_BG, line=DIVIDER, radius=0.08)
        add_round_rect(s, x + (card_w - Inches(0.6)) / 2, cy - Inches(0.3),
                       Inches(0.6), Inches(0.6),
                       ACCENT, radius=0.5)
        add_text(s, x + (card_w - Inches(0.6)) / 2, cy - Inches(0.3),
                 Inches(0.6), Inches(0.6),
                 f"0{i+1}", size=14, bold=True, color=DARK_TEXT,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x, cy + Inches(0.55), card_w, Inches(0.4),
                 title, size=14, bold=True, color=PRIMARY,
                 align=PP_ALIGN.CENTER)
        add_text(s, x + Inches(0.15), cy + Inches(1.0),
                 card_w - Inches(0.3), Inches(1.25),
                 body, size=10, color=BODY_TEXT,
                 align=PP_ALIGN.CENTER, spacing=1.3)
        if i < n - 1:
            ax1 = x + card_w + Inches(0.04)
            ax2 = x + card_w + gap - Inches(0.04)
            add_line(s, ax1, cy + Inches(1.1),
                     ax2, cy + Inches(1.1), ACCENT, 1.5)

    # bottom flow box
    add_round_rect(s, Inches(0.55), Inches(5.05), Inches(12.2), Inches(2.0),
                   SOFT_BG, radius=0.06)
    add_text(s, Inches(0.85), Inches(5.15), Inches(11.6), Inches(0.4),
             "Under the hood — what's moving where",
             size=13, bold=True, color=PRIMARY)
    add_bullets(s, Inches(0.85), Inches(5.55), Inches(11.6), Inches(1.5), [
        "Browser  →  Streamlit (Python)  →  PyMySQL  →  MySQL "
        "(users, doctors, slots, visits, prescriptions).",
        "Two AI side-paths: LangChain + Groq for the in-app care copilot, "
        "and a Gemini/Groq voice widget for always-on help.",
        "Outbound: Gmail SMTP delivers doctor credentials; "
        "ReportLab + pdfplumber export medical records.",
    ], size=11)


# ────────────────────────────────────────────────────────────────────
# SLIDE 4 — Components (3 roles)
# ────────────────────────────────────────────────────────────────────
def slide_4():
    s = prs.slides.add_slide(BLANK)
    page_bg(s)
    add_title(s, "Section 03", "Core components",
             "Three role-aware dashboards, all reading from one MySQL DB.")

    cols = [
        ("Patient", "pages/patient_dashboard.py", PRIMARY, [
            "Profile + medical history view",
            "Browse doctors by specialty",
            "Book / reschedule / cancel slots",
            "Symptom checker (AI)",
            "Prescription & PDF downloads",
            "Floating voice chatbot 24×7",
        ]),
        ("Doctor", "pages/doctor_dashboard.py", SECONDARY, [
            "Today's schedule + queue",
            "AI pre-summary of patient",
            "Diagnosis, vitals, notes",
            "Write / edit prescriptions",
            "Manage availability slots",
            "Password reset via email",
        ]),
        ("Admin", "pages/admin_dashboard.py", PURPLE, [
            "Onboard doctors (email)",
            "Activate / deactivate users",
            "Specialty & slot oversight",
            "Analytics: visits, revenue",
            "Reset doctor passwords",
            "Full DB-backed audit view",
        ]),
    ]
    col_w = Inches(4.05)
    gap   = Inches(0.15)
    start_x = Inches(0.55)
    cy = Inches(2.0)
    ch = Inches(4.95)

    for i, (role, file, col, items) in enumerate(cols):
        x = start_x + (col_w + gap) * i
        add_round_rect(s, x, cy, col_w, ch, CARD_BG, line=DIVIDER, radius=0.05)
        # header strip — taller, with room for label + icon + filename
        add_rect(s, x, cy, col_w, Inches(0.85), col)
        # role label
        add_text(s, x + Inches(0.3), cy + Inches(0.08),
                 col_w - Inches(1.2), Inches(0.4),
                 role, size=18, bold=True, color=WHITE)
        # filename under role
        add_text(s, x + Inches(0.3), cy + Inches(0.48),
                 col_w - Inches(1.2), Inches(0.32),
                 file, size=10, italic=True, color=ACCENT_SOFT)
        # round icon on right of header
        add_round_rect(s, x + col_w - Inches(0.75),
                       cy + Inches(0.15),
                       Inches(0.55), Inches(0.55),
                       WHITE, radius=0.5)
        add_text(s, x + col_w - Inches(0.75),
                 cy + Inches(0.15),
                 Inches(0.55), Inches(0.55),
                 role[0], size=20, bold=True, color=col,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        # bullets — tighter, fit the inner card
        add_bullets(s, x + Inches(0.25), cy + Inches(1.0),
                    col_w - Inches(0.5), ch - Inches(1.2),
                    items, size=10.5, spacing=1.2)

        # bottom strip
        add_rect(s, x, cy + ch - Inches(0.1), col_w, Inches(0.1), col)


# ────────────────────────────────────────────────────────────────────
# SLIDE 5 — AI implementations
# ────────────────────────────────────────────────────────────────────
def slide_5():
    s = prs.slides.add_slide(BLANK)
    page_bg(s)
    add_title(s, "Section 04", "AI implementations",
             "Four places where an LLM actually does something useful.")

    cards = [
        {"tag": "01  •  IN-APP COPILOT",
         "title": "AI Care",
         "stack": "LangChain + Groq (openai/gpt-oss-120b, fallback llama-3.1-8b-instant)",
         "what": "Database-aware assistant. Reads users, doctors, slots and visits "
                 "tables to answer 'who's free tomorrow?', 'show my last visit', "
                 "'book Dr. Sharma at 4pm'. Also surfaces patient history to "
                 "the doctor before a consult.",
         "file": "ai_care.py",
         "color": PRIMARY},
        {"tag": "02  •  ALWAYS-ON HELPER",
         "title": "Floating Voice Chatbot",
         "stack": "Gemini API + Groq fallback  •  Web Speech (STT + TTS)  •  Glass UI",
         "what": "A mic-button widget pinned bottom-right on every page. Patients "
                 "speak their question, the JS mic captures it, the model answers, "
                 "and the response is read back aloud via the browser.",
         "file": "voice_chatbot.py",
         "color": ACCENT},
        {"tag": "03  •  PUBLIC (PRE-LOGIN)",
         "title": "Public LLM Chat",
         "stack": "ChatGroq  •  LangChain Core",
         "what": "A small 'ask anything' panel on the home page for visitors who "
                 "haven't signed in yet. Answers FAQs about the product and "
                 "gives a feel for what the AI does inside the app.",
         "file": "chatbot.py",
         "color": SECONDARY},
        {"tag": "04  •  PRE-VISIT BRIEF",
         "title": "Doctor's AI Summary",
         "stack": "Prompt-engineered on Groq  •  reads DB context per patient",
         "what": "Before each consultation, the doctor gets a 5-bullet brief: "
                 "key history, recent vitals, current meds, red flags, "
                 "questions to ask — generated from the same DB the app uses.",
         "file": "ai_care.py (consultation mode)",
         "color": PURPLE},
    ]
    # 2x2 grid — bigger cards, plenty of inner padding
    cw = Inches(6.05); ch = Inches(2.5)
    gx = Inches(0.15); gy = Inches(0.15)
    x0 = Inches(0.55); y0 = Inches(2.0)

    for i, c in enumerate(cards):
        r, col = i // 2, i % 2
        x = x0 + (cw + gx) * col
        y = y0 + (ch + gy) * r

        add_round_rect(s, x, y, cw, ch, CARD_BG, line=DIVIDER, radius=0.05)
        add_rect(s, x, y, Inches(0.16), ch, c["color"])

        add_text(s, x + Inches(0.4), y + Inches(0.1),
                 cw - Inches(0.6), Inches(0.3),
                 c["tag"], size=9.5, bold=True, color=c["color"])
        add_text(s, x + Inches(0.4), y + Inches(0.38),
                 cw - Inches(0.6), Inches(0.45),
                 c["title"], size=17, bold=True, color=DARK_TEXT)
        add_text(s, x + Inches(0.4), y + Inches(0.83),
                 cw - Inches(0.6), Inches(0.32),
                 c["stack"], size=10, italic=True, color=MUTED_TEXT)
        add_text(s, x + Inches(0.4), y + Inches(1.18),
                 cw - Inches(0.6), Inches(1.0),
                 c["what"], size=10.5, color=BODY_TEXT, spacing=1.3)
        # file footer line
        add_line(s, x + Inches(0.4), y + ch - Inches(0.42),
                 x + cw - Inches(0.2), y + ch - Inches(0.42),
                 DIVIDER, 0.5)
        add_text(s, x + Inches(0.4), y + ch - Inches(0.35),
                 cw - Inches(0.6), Inches(0.3),
                 f"file: {c['file']}", size=9.5, italic=True,
                 color=MUTED_TEXT)


# ────────────────────────────────────────────────────────────────────
# SLIDE 6 — Tool stack (overview)
# ────────────────────────────────────────────────────────────────────
def slide_6():
    s = prs.slides.add_slide(BLANK)
    page_bg(s)
    add_title(s, "Section 05", "Tool stack — at a glance",
             "Four families. Full library list with working notes → next slide.")

    # 2x2 grid of category overview cards
    cats = [
        {"head": "Frontend / App shell",
         "color": PRIMARY,
         "libs": "Streamlit + custom CSS",
         "tag": "Powers the whole UI: pages, routing, forms, session state."},
        {"head": "Backend / Data",
         "color": ACCENT,
         "libs": "Python + PyMySQL + MySQL + SQLite + dotenv",
         "tag": "Server logic, schemas, dialect-aware CRUD, secrets."},
        {"head": "AI / LLM",
         "color": SECONDARY,
         "libs": "LangChain + Groq + Gemini",
         "tag": "Two chat brains: in-app copilot and floating voice widget."},
        {"head": "Utilities & I/O",
         "color": PURPLE,
         "libs": "Pillow, ReportLab, pdfplumber, pandas, requests, smtplib",
         "tag": "PDFs, charts, OAuth, SMTP — the boring plumbing that ships."},
    ]
    cw = Inches(6.05); ch = Inches(2.45)
    gx = Inches(0.15); gy = Inches(0.15)
    x0 = Inches(0.55); y0 = Inches(2.0)

    for i, cat in enumerate(cats):
        r, col = i // 2, i % 2
        x = x0 + (cw + gx) * col
        y = y0 + (ch + gy) * r
        add_round_rect(s, x, y, cw, ch, CARD_BG, line=DIVIDER, radius=0.04)
        add_rect(s, x, y, cw, Inches(0.6), cat["color"])
        add_text(s, x + Inches(0.3), y, cw - Inches(0.6), Inches(0.6),
                 cat["head"], size=14, bold=True, color=WHITE,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + Inches(0.3), y + Inches(0.8),
                 cw - Inches(0.6), Inches(0.4),
                 "LIBRARIES", size=9, bold=True, color=cat["color"])
        add_text(s, x + Inches(0.3), y + Inches(1.1),
                 cw - Inches(0.6), Inches(0.5),
                 cat["libs"], size=12, bold=True, color=DARK_TEXT,
                 spacing=1.15)
        add_line(s, x + Inches(0.3), y + Inches(1.7),
                 x + cw - Inches(0.3), y + Inches(1.7),
                 DIVIDER, 0.5)
        add_text(s, x + Inches(0.3), y + Inches(1.78),
                 cw - Inches(0.6), Inches(0.6),
                 cat["tag"], size=10.5, italic=True,
                 color=BODY_TEXT, spacing=1.3)


# ────────────────────────────────────────────────────────────────────
# SLIDE 7 — Libraries — AI / LLM  (the new full-detail slide)
# ────────────────────────────────────────────────────────────────────
def slide_7_ai_libs():
    s = prs.slides.add_slide(BLANK)
    page_bg(s)
    add_title(s, "Section 06 • Part 1", "Libraries — AI / LLM",
             "What each AI package does, and where it's used in the code.")

    # Compact table layout — header + 7 rows
    rows = [
        ("streamlit",          "≥1.30",  "UI shell",
         "Renders the whole web app — pages, widgets, session state, file uploads. "
         "We use it for every dashboard and the public chatbot."),
        ("langchain-core",     "≥0.2",   "AI plumbing",
         "Message types (HumanMessage), prompt templates and runnable interfaces. "
         "Used by ai_care.py and chatbot.py to wrap LLM calls."),
        ("langchain-groq",     "≥0.1",   "AI brain (primary)",
         "ChatGroq wrapper around Groq's hosted LLMs. Default model "
         "openai/gpt-oss-120b; falls back to llama-3.1-8b-instant for speed."),
        ("groq",               "≥0.37",  "AI brain (direct)",
         "Direct Groq SDK used as a secondary fallback path when LangChain's "
         "wrapper hits errors — short, latency-sensitive calls."),
        ("Gemini API",         "—",      "Voice chatbot brain",
         "Google's Gemini model is called from the browser-side JS in "
         "voice_chatbot.py. Used because it handles conversational health Q&A well."),
        ("Pillow (PIL)",       "≥10",    "Images",
         "Loads the logo JPG at runtime for the favicon, exposes it to "
         "st.set_page_config."),
        ("reportlab",          "≥4",     "PDF generation",
         "Generates the medical-record PDF the patient downloads after a visit "
         "(prescription, vitals, doctor notes)."),
        ("pdfplumber",         "≥0.10",  "PDF parsing",
         "Reads PDF uploads (e.g. lab reports) so the doctor can attach parsed "
         "text to the patient's record."),
    ]

    # Header row
    x0 = Inches(0.55); y0 = Inches(2.0)
    cw = (Inches(12.2))
    col_w = [Inches(2.0), Inches(1.0), Inches(2.2), Inches(7.0)]
    headers = ["Library", "Version", "Role", "What it does / where it's used"]
    x = x0
    for i, h in enumerate(headers):
        add_rect(s, x, y0, col_w[i], Inches(0.45), PRIMARY)
        add_text(s, x + Inches(0.1), y0, col_w[i] - Inches(0.2), Inches(0.45),
                 h, size=11, bold=True, color=WHITE,
                 anchor=MSO_ANCHOR.MIDDLE)
        x += col_w[i]

    # Data rows
    row_h = Inches(0.55)
    yy = y0 + Inches(0.45)
    for idx, (lib, ver, role, what) in enumerate(rows):
        bg = CARD_BG if idx % 2 == 0 else LIGHT_BG
        x = x0
        cells = [lib, ver, role, what]
        align_map = [PP_ALIGN.LEFT, PP_ALIGN.LEFT, PP_ALIGN.LEFT, PP_ALIGN.LEFT]
        bold_map  = [True, False, True, False]
        size_map  = [10.5, 10, 10, 10]
        for i, txt in enumerate(cells):
            add_rect(s, x, yy, col_w[i], row_h, bg, line=DIVIDER)
            add_text(s, x + Inches(0.1), yy,
                     col_w[i] - Inches(0.2), row_h,
                     txt, size=size_map[i], bold=bold_map[i],
                     color=DARK_TEXT if bold_map[i] else BODY_TEXT,
                     align=align_map[i], anchor=MSO_ANCHOR.MIDDLE,
                     spacing=1.15)
            x += col_w[i]
        yy += row_h


# ────────────────────────────────────────────────────────────────────
# SLIDE 8 — Libraries — Backend, Frontend, Utilities
# ────────────────────────────────────────────────────────────────────
def slide_8_other_libs():
    s = prs.slides.add_slide(BLANK)
    page_bg(s)
    add_title(s, "Section 06 • Part 2", "Libraries — Backend, Frontend, Utilities",
             "The non-AI half. Some are pip-installed, some are Python stdlib.")

    rows = [
        ("pymysql",       "≥1.1",     "Backend",
         "MySQL driver. Opens connections from db.py, runs schema CREATE statements, "
         "and executes every CRUD query."),
        ("cryptography",  "≥41",      "Backend",
         "Required transitively by PyMySQL for the secure MySQL auth handshake "
         "(caching_sha2_password)."),
        ("python-dotenv", "≥1.0",     "Backend",
         "Loads .env into os.environ so secrets (GROQ_API_KEY, SMTP creds, OAuth "
         "client) never sit in source."),
        ("sqlite3 (stdlib)", "—",     "Backend",
         "Fallback DB used when MySQL isn't reachable — same dialect-aware "
         "schema builder in db.py switches AUTOINCREMENT vs AUTO_INCREMENT."),
        ("hashlib (stdlib)", "—",     "Backend",
         "SHA-256 password hashing in auth.py. Used instead of bcrypt so the "
         "demo can run with zero install beyond Streamlit."),
        ("smtplib (stdlib)", "—",     "Backend",
         "Sends the HTML welcome email to newly onboarded doctors via Gmail."),
        ("streamlit",     "≥1.30",    "Frontend",
         "Same library listed in the AI slide — every visible page is a Streamlit "
         "script. Custom CSS is injected via st.markdown for theming."),
        ("Pillow",        "≥10",      "Frontend",
         "Logo + favicon rendering (also listed above for the AI slide context)."),
        ("pandas",        "≥2.0",     "Utility",
         "Admin analytics: visit counts, revenue, doctor load — all rendered "
         "from a pandas DataFrame."),
        ("requests",      "≥2.31",    "Utility",
         "Google OAuth flow — token exchange + userinfo call. Also used for any "
         "external HTTP call where needed."),
        ("Custom CSS",    "—",        "Frontend",
         "Hand-rolled in app.py and pages/shared_styles.py — Inter / Plus Jakarta "
         "Sans, glassmorphism cards, medical-slate palette."),
        ("HTML/JS in components.v1", "—", "Frontend",
         "The voice chatbot's mic + STT + TTS widget is built as a self-contained "
         "HTML/JS blob rendered via streamlit.components.v1.html()."),
    ]

    x0 = Inches(0.55); y0 = Inches(2.0)
    col_w = [Inches(2.4), Inches(0.9), Inches(1.7), Inches(7.2)]
    headers = ["Library", "Version", "Layer", "What it does / where it's used"]
    x = x0
    for i, h in enumerate(headers):
        add_rect(s, x, y0, col_w[i], Inches(0.42), SECONDARY)
        add_text(s, x + Inches(0.1), y0, col_w[i] - Inches(0.2), Inches(0.42),
                 h, size=10.5, bold=True, color=WHITE,
                 anchor=MSO_ANCHOR.MIDDLE)
        x += col_w[i]

    row_h = Inches(0.395)
    yy = y0 + Inches(0.42)
    for idx, (lib, ver, role, what) in enumerate(rows):
        bg = CARD_BG if idx % 2 == 0 else LIGHT_BG
        x = x0
        cells = [lib, ver, role, what]
        bold_map  = [True, False, True, False]
        size_map  = [9.5, 9, 9, 9]
        for i, txt in enumerate(cells):
            add_rect(s, x, yy, col_w[i], row_h, bg, line=DIVIDER)
            add_text(s, x + Inches(0.1), yy,
                     col_w[i] - Inches(0.2), row_h,
                     txt, size=size_map[i], bold=bold_map[i],
                     color=DARK_TEXT if bold_map[i] else BODY_TEXT,
                     align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE,
                     spacing=1.1)
            x += col_w[i]
        yy += row_h


# ────────────────────────────────────────────────────────────────────
# SLIDE 9 — Frontend details
# ────────────────────────────────────────────────────────────────────
def slide_9():
    s = prs.slides.add_slide(BLANK)
    page_bg(s)
    add_title(s, "Section 07", "Frontend — what the user sees",
             "Streamlit pages, with hand-rolled styling on top.")

    # left — page list
    add_text(s, Inches(0.55), Inches(2.0), Inches(6.2), Inches(0.4),
             "Pages in the app", size=14, bold=True, color=PRIMARY)
    pages = [
        ("app.py",                       "Public landing — login, signup, Google OAuth, public chatbot."),
        ("auth.py",                      "Hashing, seed users, OAuth handlers."),
        ("pages/patient_dashboard.py",   "Patient UI — history, booking, AI."),
        ("pages/doctor_dashboard.py",    "Doctor UI — schedule, notes, AI."),
        ("pages/admin_dashboard.py",     "Admin UI — onboard, analytics, reset."),
        ("pages/shared_styles.py",       "Common CSS + theme tokens."),
    ]
    y = Inches(2.5)
    for path, desc in pages:
        add_round_rect(s, Inches(0.55), y, Inches(6.2), Inches(0.65),
                       CARD_BG, line=DIVIDER, radius=0.16)
        add_rect(s, Inches(0.55), y, Inches(0.12), Inches(0.65), ACCENT)
        add_text(s, Inches(0.85), y, Inches(2.4), Inches(0.65),
                 path, size=11, bold=True, color=DARK_TEXT,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(3.25), y, Inches(3.4), Inches(0.65),
                 desc, size=10, color=MUTED_TEXT, italic=True,
                 anchor=MSO_ANCHOR.MIDDLE, spacing=1.2)
        y += Inches(0.72)

    # right — design system
    add_text(s, Inches(7.0), Inches(2.0), Inches(5.9), Inches(0.4),
             "Design system (hand-rolled, in CSS)",
             size=14, bold=True, color=PRIMARY)
    design = [
        ("Theme",       "Medical white + glassmorphism overlays."),
        ("Typography",  "Inter + Plus Jakarta Sans via Google Fonts."),
        ("Palette",     "Slate text, sky/teal accents, soft slate backgrounds."),
        ("Backgrounds", "Base64-encoded image + soft overlay per page."),
        ("Components",  "Cards, chips, pill buttons, sticky header, "
                        "floating voice button."),
        ("Responsive",  "Wide layout by default; cards collapse on narrow."),
    ]
    y = Inches(2.5)
    for k, v in design:
        add_round_rect(s, Inches(7.0), y, Inches(5.9), Inches(0.65),
                       SOFT_BG, radius=0.16)
        add_text(s, Inches(7.2), y, Inches(1.6), Inches(0.65),
                 k, size=11, bold=True, color=PRIMARY,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(8.85), y, Inches(3.95), Inches(0.65),
                 v, size=10, color=BODY_TEXT,
                 anchor=MSO_ANCHOR.MIDDLE, spacing=1.2)
        y += Inches(0.72)


# ────────────────────────────────────────────────────────────────────
# SLIDE 10 — Backend details
# ────────────────────────────────────────────────────────────────────
def slide_10():
    s = prs.slides.add_slide(BLANK)
    page_bg(s)
    add_title(s, "Section 08", "Backend — what runs server-side",
             "Pure Python. MySQL at the center, two AI side-paths, SMTP for mail.")

    # Center MySQL box
    add_round_rect(s, Inches(4.9), Inches(2.05), Inches(3.5), Inches(1.15),
                   PRIMARY, radius=0.08)
    add_text(s, Inches(4.9), Inches(2.05), Inches(3.5), Inches(1.15),
             "MySQL (SQLite fallback)\nusers • doctors • slots\nvisits • prescriptions",
             size=12, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, spacing=1.2)

    # 4 satellites
    sats = [
        ("auth.py",         "Hash, OAuth,\nseed users",        SECONDARY, Inches(0.6),  Inches(3.85)),
        ("db.py",           "Schema, CRUD,\ndynamic config",  ACCENT,    Inches(8.95), Inches(3.85)),
        ("email_service.py","SMTP + HTML\ndoctor welcome",   PINK,      Inches(0.6),  Inches(5.5)),
        ("ai_care.py",      "LangChain + Groq\nDB-aware AI",  PURPLE,    Inches(8.95), Inches(5.5)),
    ]
    for name, body, col, x, y in sats:
        add_round_rect(s, x, y, Inches(3.5), Inches(1.15),
                       CARD_BG, line=col, radius=0.06)
        add_text(s, x, y + Inches(0.1), Inches(3.5), Inches(0.4),
                 name, size=12, bold=True, color=col, align=PP_ALIGN.CENTER)
        add_text(s, x, y + Inches(0.55), Inches(3.5), Inches(0.55),
                 body, size=10, color=BODY_TEXT,
                 align=PP_ALIGN.CENTER, spacing=1.2)

    # connectors (straight, calculated endpoints)
    add_line(s, Inches(2.85), Inches(4.4),  Inches(4.9),  Inches(2.85), ACCENT, 1.5)
    add_line(s, Inches(8.4),  Inches(2.85), Inches(10.45), Inches(4.4),  ACCENT, 1.5)
    add_line(s, Inches(2.85), Inches(6.05), Inches(4.9),   Inches(3.2),  ACCENT, 1.5)
    add_line(s, Inches(8.4),  Inches(3.2),  Inches(10.45), Inches(6.05), ACCENT, 1.5)


# ────────────────────────────────────────────────────────────────────
# SLIDE 11 — Features (the long list)
# ────────────────────────────────────────────────────────────────────
def slide_11():
    s = prs.slides.add_slide(BLANK)
    page_bg(s)
    add_title(s, "Section 09", "Features — the long list",
             "Twelve things that actually work today, not just in spec.")

    features = [
        ("Role-aware auth",       "Email+password + Google OAuth, session state per role."),
        ("Smart booking",         "Patients see only valid future slots; doctors see only their queue."),
        ("AI symptom checker",    "Patient describes symptoms → AI returns structured triage."),
        ("Doctor AI brief",       "Auto-generated 5-bullet patient summary before each consult."),
        ("Prescription writer",   "Add medicine + dosage + duration; visible to patient."),
        ("PDF medical record",    "ReportLab-generated downloadable record."),
        ("Voice chatbot",         "Floating mic button on every page; speaks back via Web Speech."),
        ("Voice / text fallback", "Mic-less users still get text chat, same AI brain."),
        ("Doctor onboarding",     "Admin creates doctor → email goes out with temp password."),
        ("Password reset",        "Admin can reset any doctor's password."),
        ("Activate / deactivate", "Admin can suspend users in one click."),
        ("Analytics dashboard",   "Visit counts, revenue, doctor load — pandas-powered."),
    ]
    cw = Inches(3.0); ch = Inches(1.55)
    gx = Inches(0.12); gy = Inches(0.15)
    x0 = Inches(0.55); y0 = Inches(2.0)

    palette = [PRIMARY, SECONDARY, ACCENT, PURPLE]
    for i, (title, body) in enumerate(features):
        r, col = i // 4, i % 4
        x = x0 + (cw + gx) * col
        y = y0 + (ch + gy) * r
        col_color = palette[col]
        add_round_rect(s, x, y, cw, ch, CARD_BG, line=DIVIDER, radius=0.08)
        add_rect(s, x, y, Inches(0.1), ch, col_color)
        add_text(s, x + Inches(0.25), y + Inches(0.1),
                 Inches(0.7), Inches(0.3),
                 f"#{i+1:02d}", size=10, bold=True, color=col_color)
        add_text(s, x + Inches(0.25), y + Inches(0.4),
                 cw - Inches(0.4), Inches(0.4),
                 title, size=12, bold=True, color=DARK_TEXT)
        add_text(s, x + Inches(0.25), y + Inches(0.78),
                 cw - Inches(0.4), Inches(0.7),
                 body, size=9.5, color=MUTED_TEXT,
                 italic=True, spacing=1.25)


# ────────────────────────────────────────────────────────────────────
# SLIDE 12 — Core concepts
# ────────────────────────────────────────────────────────────────────
def slide_12():
    s = prs.slides.add_slide(BLANK)
    page_bg(s)
    add_title(s, "Section 10", "Core concepts",
             "The handful of ideas everything else hangs off.")

    concepts = [
        ("Role", "The single switch that decides which dashboard renders. "
                 "Stored on the users table, enforced server-side, never trusted from the browser."),
        ("Visit", "An appointment becomes a visit only after the doctor marks "
                  "it complete — that's when AI summary, prescription, and PDF kick in."),
        ("Slot", "A (doctor, weekday, time-window) tuple. Free / booked is "
                 "computed, not stored — fewer race conditions."),
        ("AI context", "Every LLM call is given the same DB slice the user is "
                       "looking at. This is what makes 'book Dr. Sharma at 4pm' actually work."),
        ("Voice loop", "Browser mic → STT → backend LLM → TTS → speaker. "
                       "Each step falls back independently if the next one fails."),
        ("One DB, three views", "Patients, doctors and admins all read the "
                                "same rows — only the columns and actions differ."),
    ]
    cw = Inches(6.05); ch = Inches(1.55)
    gx = Inches(0.15); gy = Inches(0.15)
    x0 = Inches(0.55); y0 = Inches(2.0)
    palette = [PRIMARY, ACCENT, SECONDARY, PURPLE, PINK, WARN]

    for i, (head, body) in enumerate(concepts):
        r, col = i // 2, i % 2
        x = x0 + (cw + gx) * col
        y = y0 + (ch + gy) * r
        c = palette[i]
        add_round_rect(s, x, y, cw, ch, CARD_BG, line=DIVIDER, radius=0.05)
        add_round_rect(s, x + Inches(0.25), y + Inches(0.25),
                       Inches(0.55), Inches(0.55),
                       c, radius=0.5)
        add_text(s, x + Inches(0.25), y + Inches(0.25),
                 Inches(0.55), Inches(0.55),
                 str(i+1), size=16, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + Inches(1.0), y + Inches(0.18),
                 cw - Inches(1.2), Inches(0.4),
                 head, size=13.5, bold=True, color=DARK_TEXT)
        add_text(s, x + Inches(1.0), y + Inches(0.58),
                 cw - Inches(1.2), Inches(0.95),
                 body, size=10.5, color=BODY_TEXT, spacing=1.3)


# ────────────────────────────────────────────────────────────────────
# SLIDE 13 — Difficulties
# ────────────────────────────────────────────────────────────────────
def slide_13():
    s = prs.slides.add_slide(BLANK)
    page_bg(s)
    add_title(s, "Section 11", "Difficulties, honestly",
             "Things that took longer than they should have.")

    items = [
        ("MySQL not running locally",
         "Most laptops don't have MySQL. First launch crashed. "
         "Fix: dynamic db_config + UI prompt to enter creds."),
        ("Streamlit reruns everything",
         "Any click reruns the whole script; LLM calls fired multiple times. "
         "Fix: st.session_state for chat history."),
        ("OAuth redirect in dev",
         "Google OAuth needs exact redirect URIs. localhost:8501 worked, "
         "127.0.0.1 didn't. Hard-coded the redirect URI in auth.py."),
        ("AI cost + latency",
         "Long LLM calls blocked the UI. Fix: st.spinner, smaller prompts, "
         "llama-3.1-8b-instant fallback for short queries."),
        ("Voice in the browser",
         "Web Speech needs user gesture + HTTPS (or localhost). "
         "Worked in dev, will need HTTPS in production."),
        ("Schema drift MySQL/SQLite",
         "AUTO_INCREMENT vs AUTOINCREMENT, BOOL vs BOOLEAN. "
         "Fix: dialect-aware schema builder in db.py."),
    ]
    cw = Inches(6.05); ch = Inches(1.6)
    gx = Inches(0.15); gy = Inches(0.13)
    x0 = Inches(0.55); y0 = Inches(2.0)

    for i, (title, body) in enumerate(items):
        r, col = i // 2, i % 2
        x = x0 + (cw + gx) * col
        y = y0 + (ch + gy) * r
        add_round_rect(s, x, y, cw, ch, CARD_BG, line=DIVIDER, radius=0.04)
        add_rect(s, x, y, Inches(0.12), ch, WARN)
        add_text(s, x + Inches(0.3), y + Inches(0.1),
                 cw - Inches(0.5), Inches(0.35),
                 title, size=11.5, bold=True, color=DARK_TEXT)
        add_text(s, x + Inches(0.3), y + Inches(0.46),
                 cw - Inches(0.5), Inches(1.1),
                 body, size=9.5, color=BODY_TEXT, italic=True, spacing=1.25)


# ────────────────────────────────────────────────────────────────────
# SLIDE 14 — Closing
# ────────────────────────────────────────────────────────────────────
def slide_14():
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SW, SH, PAGE_BG)

    add_rect(s, 0, 0, SW, Inches(3.5), PRIMARY)
    add_rect(s, 0, Inches(3.5), SW, Inches(0.06), ACCENT)
    c1 = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.5), Inches(-1.5),
                            Inches(3.5), Inches(3.5))
    c1.fill.solid(); c1.fill.fore_color.rgb = SECONDARY
    c1.line.fill.background(); c1.shadow.inherit = False
    c2 = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(11.8), Inches(0.8),
                            Inches(2.0), Inches(2.0))
    c2.fill.solid(); c2.fill.fore_color.rgb = ACCENT
    c2.line.fill.background(); c2.shadow.inherit = False

    add_text(s, Inches(0.7), Inches(0.7), Inches(11), Inches(0.4),
             "THANK YOU", size=14, bold=True, color=ACCENT)
    add_text(s, Inches(0.7), Inches(1.05), Inches(11), Inches(1.2),
             "That's PCMHS.", size=44, bold=True, color=WHITE)
    add_text(s, Inches(0.7), Inches(2.2), Inches(10), Inches(0.9),
             "A patient-care app that doesn't pretend AI is magic — "
             "it just moves the right information to the right person, faster.",
             size=14, italic=True, color=ACCENT_SOFT, spacing=1.4)

    cards = [
        ("What you'd build next", PRIMARY, [
            "Multi-language chatbot (Hindi, Spanish).",
            "Doctor-side calendar sync (Google Calendar).",
            "Insurance + billing module.",
            "Mobile-first layout (current is desktop-tuned)."]),
        ("What I'd watch", WARN, [
            "LLM cost scaling with concurrent doctors.",
            "MySQL connection pool under load.",
            "PII safety in AI prompts — no raw history leaks.",
            "Voice latency on low-end devices."]),
        ("Where to look first", ACCENT, [
            "app.py — the entry point & theme.",
            "pages/ai_care.py — the AI brain.",
            "db.py — schema and dialect switching.",
            ".env.example — keys to bring your own."]),
    ]
    cw = Inches(4.05); ch = Inches(3.05)
    gx = Inches(0.15)
    x0 = Inches(0.55); y = Inches(3.85)

    for i, (head, col, items) in enumerate(cards):
        x = x0 + (cw + gx) * i
        add_round_rect(s, x, y, cw, ch, CARD_BG, line=DIVIDER, radius=0.04)
        add_rect(s, x, y, cw, Inches(0.55), col)
        add_text(s, x + Inches(0.3), y, cw - Inches(0.6), Inches(0.55),
                 head, size=13, bold=True, color=WHITE,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_bullets(s, x + Inches(0.3), y + Inches(0.75),
                    cw - Inches(0.5), ch - Inches(0.9),
                    items, size=10.5, spacing=1.4)


# ────────────────────────────────────────────────────────────────────
# Build & save
# ────────────────────────────────────────────────────────────────────
slides = [
    slide_1, slide_2, slide_3, slide_4, slide_5, slide_6,
    slide_7_ai_libs, slide_8_other_libs, slide_9, slide_10,
    slide_11, slide_12, slide_13, slide_14,
]
total = len(slides)
for i, fn in enumerate(slides, 1):
    fn()
    add_footer(prs.slides[-1], i, total)

out_path = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                        "PCMHS_Project_Presentation.pptx")
prs.save(out_path)
print(f"[OK] Saved: {out_path}")
print(f"     Slides: {total}")
