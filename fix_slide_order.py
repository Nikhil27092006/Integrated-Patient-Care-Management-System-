"""Reorder IPCMS_v2.pptx so THANK YOU is last and screenshots are in order.

Current bad order:
  19: SCREENSHOT 01
  20: SCREENSHOT 05
  21: SCREENSHOT 03
  22: THANK YOU
  23: SCREENSHOT 02
  24: SCREENSHOT 04

Desired order:
  19: SCREENSHOT 01
  20: SCREENSHOT 02
  21: SCREENSHOT 03
  22: SCREENSHOT 04
  23: SCREENSHOT 05
  24: THANK YOU
"""

import re
from pptx import Presentation
from pptx.oxml.ns import qn

SRC = "IPCMS_v2.pptx"


def slide_first_text(slide):
    for s in slide.shapes:
        if s.has_text_frame and s.text_frame.text.strip():
            return s.text_frame.text.strip()
    return ""


def move_slide(prs, old_index, new_index):
    sldIdLst = prs.slides._sldIdLst
    slides = list(sldIdLst)
    sldIdLst.remove(slides[old_index])
    sldIdLst.insert(new_index, slides[old_index])


def main():
    prs = Presentation(SRC)

    # Map: SCREENSHOT N -> current index, THANK YOU -> current index
    current = {}
    for i, slide in enumerate(prs.slides):
        first = slide_first_text(slide).split("\n")[0].strip()
        m = re.match(r"^SCREENSHOT\s+(\d+)$", first)
        if m:
            current[("shot", int(m.group(1)))] = i
        elif first == "THANK YOU":
            current[("thanks",)] = i

    print("Current positions:")
    print(f"  SCREENSHOT 01 -> {current[('shot', 1)]}")
    print(f"  SCREENSHOT 02 -> {current[('shot', 2)]}")
    print(f"  SCREENSHOT 03 -> {current[('shot', 3)]}")
    print(f"  SCREENSHOT 04 -> {current[('shot', 4)]}")
    print(f"  SCREENSHOT 05 -> {current[('shot', 5)]}")
    print(f"  THANK YOU     -> {current[('thanks',)]}")

    total = len(prs.slides)
    # Final order: SCREENSHOT 01..05 occupy indices [total-6 .. total-2],
    # THANK YOU occupies total-1.

    # Strategy: move from last to first.
    # 1) THANK YOU -> last (total-1)
    move_slide(prs, current[("thanks",)], total - 1)
    # Refresh after move
    for i, slide in enumerate(prs.slides):
        first = slide_first_text(slide).split("\n")[0].strip()
        m = re.match(r"^SCREENSHOT\s+(\d+)$", first)
        if m:
            current[("shot", int(m.group(1)))] = i
        elif first == "THANK YOU":
            current[("thanks",)] = i

    # 2) Move each SCREENSHOT in reverse order to its target slot
    # After step 1, THANK YOU is at total-1.
    # Target for SCREENSHOT 05: total-2, SCREENSHOT 04: total-3, etc.
    for n in range(5, 0, -1):
        target = total - 1 - (5 - n + 1)  # = total - 7 + n
        # Actually: SCREENSHOT 05 -> total-2, SCREENSHOT 04 -> total-3, ...,
        # SCREENSHOT 01 -> total-6
        target = total - (5 - n + 1) - 1  # n=5 -> total-6+5? Let me just compute:
        # n=5: target = total - 2
        # n=4: target = total - 3
        # n=3: target = total - 4
        # n=2: target = total - 5
        # n=1: target = total - 6
        target = (total - 1) - (5 - n + 1)
        # Simplify: target = total - 1 - (6 - n) = total - 7 + n
        target = total - 7 + n
        src_idx = current[("shot", n)]
        print(f"  Move SCREENSHOT {n:02d} from {src_idx} to {target}")
        move_slide(prs, src_idx, target)
        # Refresh positions after each move
        for i, slide in enumerate(prs.slides):
            first = slide_first_text(slide).split("\n")[0].strip()
            m = re.match(r"^SCREENSHOT\s+(\d+)$", first)
            if m:
                current[("shot", int(m.group(1)))] = i
            elif first == "THANK YOU":
                current[("thanks",)] = i

    # Refresh footer page numbers across the whole deck
    final_total = len(prs.slides)
    for idx, slide in enumerate(prs.slides):
        for shp in slide.shapes:
            if shp.has_text_frame:
                for p in shp.text_frame.paragraphs:
                    for r in p.runs:
                        if r.text and re.match(r"^\s*\d+\s*/\s*\d+\s*$", r.text):
                            left, _ = r.text.split("/")
                            r.text = f"{left.strip()} / {final_total}"

    prs.save(SRC)
    print(f"\nSaved {SRC}.")

    # Print final order
    print("\nFinal order:")
    for i, slide in enumerate(prs.slides, 1):
        title = ""
        for s in slide.shapes:
            if s.has_text_frame:
                t = s.text_frame.text.strip()
                if t:
                    title = t.split("\n")[0][:60]
                    break
        print(f"  Slide {i}: {title}")


if __name__ == "__main__":
    main()
