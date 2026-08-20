"""Append problem-statement + 4 milestone slides to IPCMS_v2.pptx.

Mirrors the deck's existing style: Calibri, accent #14B8A6, deep-blue #0E4F8A,
slate #0F172A, muted #64748B, footer rule + page-number.
"""

from copy import deepcopy
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

SRC = "IPCMS_v2.pptx"

TEAL = RGBColor(0x14, 0xB8, 0xA6)
DEEP = RGBColor(0x0E, 0x4F, 0x8A)
SLATE = RGBColor(0x0F, 0x17, 0x2A)
MUTED = RGBColor(0x64, 0x74, 0x8B)
BODY = RGBColor(0x33, 0x4E, 0x68)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG_TINT = RGBColor(0xF1, 0xF5, 0xF9)
CARD_FILL = RGBColor(0xF8, 0xFA, 0xFC)


def in_to_emu(x: float) -> int:
    return int(round(x * 914400))


def add_bg(prs, slide):
    """No-op: the slide's default background is already white in this deck.

    Earlier versions of this function drew a full-bleed white rectangle and
    tried to push it to the back via spTree.insert(2, ...). That actually
    landed it ABOVE the other shapes in z-order, blanketing the whole
    slide. The cleanest fix is to trust the slide's default white bg.
    """
    return None


def add_header(prs, slide, eyebrow, title, subtitle):
    """Eyebrow + big title + subtitle + horizontal rule."""
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    in_to_emu(0.55), in_to_emu(0.55),
                                    in_to_emu(0.12), in_to_emu(0.55))
    accent.line.fill.background()
    accent.fill.solid()
    accent.fill.fore_color.rgb = TEAL

    tb = slide.shapes.add_textbox(in_to_emu(0.78), in_to_emu(0.5),
                                   in_to_emu(12.0), in_to_emu(0.28))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = eyebrow
    r.font.name = "Calibri"
    r.font.size = Pt(10)
    r.font.bold = True
    r.font.color.rgb = TEAL

    tb = slide.shapes.add_textbox(in_to_emu(0.78), in_to_emu(0.78),
                                   in_to_emu(12.2), in_to_emu(0.55))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = "Calibri"
    r.font.size = Pt(26)
    r.font.bold = True
    r.font.color.rgb = SLATE

    tb = slide.shapes.add_textbox(in_to_emu(0.78), in_to_emu(1.35),
                                   in_to_emu(12.2), in_to_emu(0.4))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = subtitle
    r.font.name = "Calibri"
    r.font.size = Pt(12)
    r.font.color.rgb = MUTED

    rule = slide.shapes.add_connector(1, in_to_emu(0.78), in_to_emu(1.78),
                                      in_to_emu(12.6), in_to_emu(1.78))
    rule.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
    rule.line.width = Pt(0.75)


def add_footer(prs, slide, page_num, total):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, in_to_emu(7.18),
                                  prs.slide_width, in_to_emu(0.32))
    bar.line.fill.background()
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(0xF8, 0xFA, 0xFC)

    tb = slide.shapes.add_textbox(in_to_emu(0.4), in_to_emu(7.22),
                                   in_to_emu(9.0), in_to_emu(0.25))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "IPCMS — Integrated Patient Care Management System"
    r.font.name = "Calibri"
    r.font.size = Pt(9)
    r.font.color.rgb = MUTED

    tb = slide.shapes.add_textbox(in_to_emu(11.6), in_to_emu(7.22),
                                   in_to_emu(1.4), in_to_emu(0.25))
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = f"{page_num} / {total}"
    r.font.name = "Calibri"
    r.font.size = Pt(9)
    r.font.color.rgb = MUTED


def add_text(slide, left, top, width, height, text,
             *, size=11, bold=False, color=BODY, name="Calibri",
             align=None):
    tb = slide.shapes.add_textbox(in_to_emu(left), in_to_emu(top),
                                   in_to_emu(width), in_to_emu(height))
    tf = tb.text_frame
    tf.word_wrap = True
    if isinstance(text, str):
        lines = [text]
    else:
        lines = text
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        if align is not None:
            p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.name = name
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return tb


def add_card(slide, left, top, width, height, *,
             fill=CARD_FILL, line_rgb=RGBColor(0xE2, 0xE8, 0xF0)):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   in_to_emu(left), in_to_emu(top),
                                   in_to_emu(width), in_to_emu(height))
    card.fill.solid()
    card.fill.fore_color.rgb = fill
    card.line.color.rgb = line_rgb
    card.line.width = Pt(0.75)
    card.shadow.inherit = False
    # gentler radius via XML
    sp = card._element
    prstGeom = sp.find(qn("a:prstGeom"))
    if prstGeom is not None:
        avLst = prstGeom.find(qn("a:avLst"))
        if avLst is None:
            avLst = etree.SubElement(prstGeom, qn("a:avLst"))
        # adjust radius (default 0.07 is close to deck; leave as-is)
    return card


def add_number_badge(slide, left, top, width, height, label):
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    in_to_emu(left), in_to_emu(top),
                                    in_to_emu(width), in_to_emu(height))
    badge.fill.solid()
    badge.fill.fore_color.rgb = TEAL
    badge.line.fill.background()
    badge.shadow.inherit = False
    tf = badge.text_frame
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = label
    r.font.name = "Calibri"
    r.font.size = Pt(13)
    r.font.bold = True
    r.font.color.rgb = WHITE
    return badge


def add_eyebrow(slide, left, top, width, text, color=TEAL):
    tb = slide.shapes.add_textbox(in_to_emu(left), in_to_emu(top),
                                   in_to_emu(width), in_to_emu(0.3))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.name = "Calibri"
    r.font.size = Pt(11)
    r.font.bold = True
    r.font.color.rgb = color


def build_section_header_slide(prs, section_label, title):
    """Big section opener like 'SECTION 02' style."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    add_bg(prs, slide)

    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    in_to_emu(0.55), in_to_emu(2.7),
                                    in_to_emu(0.18), in_to_emu(2.1))
    accent.line.fill.background()
    accent.fill.solid()
    accent.fill.fore_color.rgb = TEAL

    add_text(slide, 0.9, 2.7, 12.0, 0.5, section_label,
             size=14, bold=True, color=TEAL)
    add_text(slide, 0.9, 3.2, 12.0, 1.1, title,
             size=40, bold=True, color=SLATE)

    rule = slide.shapes.add_connector(1, in_to_emu(0.9), in_to_emu(4.45),
                                      in_to_emu(8.0), in_to_emu(4.45))
    rule.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
    rule.line.width = Pt(0.75)
    return slide


def build_problem_solutions_slide(prs, page_num, total):
    """One slide listing 4 problems on the left and their solutions on the right."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(prs, slide)
    add_header(prs, slide, "SECTION 12",
               "Problem statements & solutions",
               "What real-world gaps the project closes — and exactly how the build answers each one.")
    add_footer(prs, slide, page_num, total)

    # Layout: left column = problem, right column = solution.
    # Title bar over each column.
    add_text(slide, 0.55, 1.95, 6.0, 0.35, "PROBLEM",
             size=11, bold=True, color=DEEP)
    add_text(slide, 6.95, 1.95, 6.0, 0.35, "OUR SOLUTION",
             size=11, bold=True, color=DEEP)

    pairs = [
        (
            "Fragmented patient records",
            "Different doctors store notes in different places — history is hard to retrieve, harder to share.",
            "Single MySQL schema with role-aware views. Patients, doctors and admins all read the same tables, so history, visits and prescriptions live in one place.",
        ),
        (
            "Manual appointment routing",
            "Reception reads forms and assigns slots by hand — slow, error-prone, no self-service.",
            "Self-service slot booking by patients, validated server-side. Doctors manage availability; the chatbot can book appointments in natural language.",
        ),
        (
            "No intelligent first response",
            "Patients Google symptoms or wait days for advice; doctors can't scale 1-to-1 chat.",
            "Two AI assistants (in-app Groq copilot + always-on voice widget) give 24×7 triage, prescription context and follow-up answers using the patient's DB history.",
        ),
        (
            "Paper-based pharmacy & onboarding",
            "Buying medicine, refilling prescriptions, and onboarding new doctors all run on paper / manual calls.",
            "Digital pharmacy module with admin CRUD on medicine images, patient-side buy/book flow, e-prescription PDF export, and doctor-onboarding with auto-emailed credentials.",
        ),
    ]

    y = 2.35
    row_h = 1.05
    for i, (head, prob, sol) in enumerate(pairs):
        # Left problem card
        add_card(slide, 0.55, y, 6.0, row_h,
                 fill=RGBColor(0xFE, 0xF3, 0xF2),
                 line_rgb=RGBColor(0xF1, 0xC0, 0xB8))
        add_number_badge(slide, 0.7, y + 0.18, 0.55, 0.55, f"P{i+1}")
        add_text(slide, 1.4, y + 0.12, 5.1, 0.35, head,
                 size=12, bold=True, color=SLATE)
        add_text(slide, 1.4, y + 0.48, 5.1, 0.6, prob,
                 size=10, color=BODY)

        # Right solution card
        add_card(slide, 6.95, y, 6.0, row_h,
                 fill=RGBColor(0xEC, 0xFD, 0xF5),
                 line_rgb=RGBColor(0x99, 0xE8, 0xD2))
        badge = add_number_badge(slide, 7.1, y + 0.18, 0.55, 0.55, f"S{i+1}")
        # recolor badge to deep blue for solutions
        badge.fill.solid()
        badge.fill.fore_color.rgb = DEEP
        add_text(slide, 7.8, y + 0.12, 5.0, 0.35, f"Fix for: {head}",
                 size=12, bold=True, color=SLATE)
        add_text(slide, 7.8, y + 0.48, 5.0, 0.6, sol,
                 size=10, color=BODY)

        y += row_h + 0.07
    return slide


def build_milestone_slide(prs, page_num, total, *, number, title, scope,
                          deliverables, files, outcome):
    """Single milestone slide, mirrors the 'SECTION 04' style."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(prs, slide)
    add_header(prs, slide,
               f"MILESTONE {number:02d}",
               title,
               scope)
    add_footer(prs, slide, page_num, total)

    # Left card: deliverables
    add_card(slide, 0.55, 2.0, 6.2, 4.6)
    add_eyebrow(slide, 0.75, 2.1, 6.0, "WHAT WAS BUILT", color=DEEP)
    for i, line in enumerate(deliverables):
        add_text(slide, 0.75, 2.55 + i * 0.4, 6.0, 0.4,
                 f"•  {line}", size=11, color=BODY)

    # Right column: files touched + outcome
    add_card(slide, 6.95, 2.0, 6.0, 2.2)
    add_eyebrow(slide, 7.15, 2.1, 5.6, "FILES / MODULES", color=DEEP)
    add_text(slide, 7.15, 2.55, 5.6, 1.6, files,
             size=11, color=BODY)

    add_card(slide, 6.95, 4.35, 6.0, 2.25,
             fill=RGBColor(0xEC, 0xFD, 0xF5),
             line_rgb=RGBColor(0x99, 0xE8, 0xD2))
    add_eyebrow(slide, 7.15, 4.45, 5.6, "OUTCOME", color=DEEP)
    add_text(slide, 7.15, 4.9, 5.6, 1.6, outcome,
             size=11, color=BODY)
    return slide


def move_slide(prs, old_index, new_index):
    """Reorder slides by XML manipulation."""
    sldIdLst = prs.slides._sldIdLst
    slides = list(sldIdLst)
    sldIdLst.remove(slides[old_index])
    sldIdLst.insert(new_index, slides[old_index])


def remove_slide(prs, index):
    """Remove a slide at the given index."""
    sldIdLst = prs.slides._sldIdLst
    slides = list(sldIdLst)
    sldId = slides[index]
    # Drop the relationship from this slide to the presentation part.
    rId = sldId.get(qn("r:id"))
    prs.part.drop_rel(rId)
    sldIdLst.remove(sldId)


def slide_first_text(slide):
    for s in slide.shapes:
        if s.has_text_frame and s.text_frame.text.strip():
            return s.text_frame.text.strip()
    return ""


def purge_existing_appends(prs):
    """Remove any pre-existing SECTION 12 / MILESTONE 0X slides (idempotency).

    Only matches slides whose first non-empty textbox equals the eyebrow
    string verbatim. This avoids ever touching the THANK YOU slide or any
    other existing content.
    """
    targets = {"SECTION 12", "MILESTONE 01", "MILESTONE 02",
               "MILESTONE 03", "MILESTONE 04"}
    # Iterate from end to start so indices stay valid.
    for idx in range(len(prs.slides) - 1, -1, -1):
        first = slide_first_text(prs.slides[idx]).split("\n")[0].strip()
        if first in targets:
            remove_slide(prs, idx)


def main():
    prs = Presentation(SRC)
    purge_existing_appends(prs)
    existing = len(prs.slides)
    # Final total after appending: existing + 1 (problems) + 4 (milestones)
    final_total = existing + 5
    # Page numbers continue from existing count.
    base = existing + 1  # first new slide's page number

    # 1) Problem statements & solutions
    build_problem_solutions_slide(prs, base, final_total)

    # 2) Milestones
    milestones = [
        dict(
            number=1,
            title="Login UI — design & polish",
            scope="Built the Streamlit login/signup screens and refined the visual system around them.",
            deliverables=[
                "Email + password sign-in and sign-up forms",
                "Google OAuth entry alongside classic auth",
                "Role detection on login (patient / doctor / admin)",
                "Glassmorphism cards, sticky header, custom CSS theme",
                "Responsive layout for narrow screens",
            ],
            files="auth.py · app.py · pages/shared_styles.py",
            outcome="A polished, themed login experience that detects role and routes each user to the right dashboard without extra clicks.",
        ),
        dict(
            number=2,
            title="Chatbot + appointment booking + DB link",
            scope="Installed the AI chatbot, taught it to book appointments, and wired it to the local database.",
            deliverables=[
                "Public chatbot on the home page (pre-login)",
                "Always-on voice chatbot widget pinned to every page",
                "NLP-based appointment booking through the chatbot",
                "Chatbot connected to local MySQL/SQLite DB",
                "Reads doctors, slots, and visits tables to answer in context",
            ],
            files="chatbot.py · voice_chatbot.py · ai_care.py · db.py",
            outcome="Patients can ask in natural language — 'book Dr. Sharma tomorrow 4pm' — and the chatbot checks real DB state before confirming the slot.",
        ),
        dict(
            number=3,
            title="Pharmacy, prescriptions & patient purchases",
            scope="Added a digital pharmacy with admin-side medicine CRUD and a patient purchase flow.",
            deliverables=[
                "Medicine catalogue with images (Medicine_Images/)",
                "Admin-side CRUD: add, edit, delete medicines",
                "Doctor writes prescriptions during consultation",
                "Auto-generated e-prescription PDF (ReportLab)",
                "Patient can browse, book and buy medicines in-app",
            ],
            files="add_medicines_admin.py · pages/admin_dashboard.py · pages/patient_dashboard.py · pages/doctor_dashboard.py",
            outcome="A closed loop: doctor prescribes → patient sees it → patient can buy the same medicines through the pharmacy, all from one DB.",
        ),
        dict(
            number=4,
            title="OCR, external sources & doctor email onboarding",
            scope="Brought external data in (OCR) and pushed notifications out (email credentials).",
            deliverables=[
                "OCR module to read uploaded medical / lab reports",
                "Connection to an external data source for richer context",
                "Admin can create a doctor login from the dashboard",
                "Generated credentials are auto-mailed to that doctor",
                "Doctor uses the emailed credentials for first sign-in",
            ],
            files="email_service.py · auth.py · pages/admin_dashboard.py · pages/doctor_dashboard.py",
            outcome="Hospitals can onboard doctors in one click, patients can upload paper reports, and the system stays the single source of truth for everyone.",
        ),
    ]

    for i, m in enumerate(milestones):
        build_milestone_slide(prs, base + 1 + i, final_total, **m)

    # Move THANK YOU slide (was last at index existing-1) to the very end.
    thank_you_idx = existing - 1  # original index of THANK YOU
    new_thank_you_idx = len(prs.slides) - 1
    if thank_you_idx != new_thank_you_idx:
        move_slide(prs, thank_you_idx, new_thank_you_idx)

    # Update the page-number text on the THANK YOU slide.
    # Find the right-most text "14 / 14"-shaped run and rewrite it.
    thanks = prs.slides[-1]
    for shp in thanks.shapes:
        if shp.has_text_frame:
            for p in shp.text_frame.paragraphs:
                for r in p.runs:
                    if "/ 14" in r.text or r.text.strip().endswith("/ 14"):
                        r.text = f"{final_total} / {final_total}"
                    if "1 / 14" in r.text or "2 / 14" in r.text or "14 / 14" in r.text:
                        r.text = f"{final_total} / {final_total}"
            # Also rewrite footer page-of-N strings throughout the deck.
    # Update every existing slide's page-of-N footer.
    total_slides = len(prs.slides)
    for idx in range(total_slides - 1):
        slide = prs.slides[idx]
        for shp in slide.shapes:
            if shp.has_text_frame:
                for p in shp.text_frame.paragraphs:
                    for r in p.runs:
                        if r.text and r.text.strip().endswith("/ 14"):
                            r.text = f"{idx + 1} / {final_total}"

    prs.save(SRC)
    print(f"Saved {SRC} with {final_total} slides.")


if __name__ == "__main__":
    main()
