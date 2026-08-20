"""Check current slide order in IPCMS_v2.pptx."""
from pptx import Presentation

prs = Presentation("IPCMS_v2.pptx")
for i, slide in enumerate(prs.slides, 1):
    title = ""
    for s in slide.shapes:
        if s.has_text_frame:
            t = s.text_frame.text.strip()
            if t:
                title = t.split("\n")[0][:60]
                break
    print(f"Slide {i}: {title}")
print(f"\nTotal slides: {len(prs.slides)}")
